"""
Adjust line spacing for all text boxes in a PowerPoint file.

Usage: python adjust_ppt_line_spacing.py INPUT_FILE OUTPUT_FILE [--spacing VALUE]

Arguments:
  INPUT_FILE    Path to input PowerPoint file (.pptx)
  OUTPUT_FILE   Path to output PowerPoint file (.pptx)
  --spacing     Line spacing value (default: 1.0)
                - 1.0 = single spacing
                - 1.5 = 1.5 line spacing
                - 2.0 = double spacing
                - Can also use Pt values like "12pt" for exact point spacing

Example:
  python adjust_ppt_line_spacing.py input.pptx output.pptx --spacing 1.2
  python adjust_ppt_line_spacing.py input.pptx output.pptx --spacing 14pt
"""

import sys
from pptx import Presentation
from pptx.util import Pt

def adjust_line_spacing(input_file, output_file, spacing=1.0):
    """
    Adjust line spacing for all text boxes in a PowerPoint presentation.

    Args:
        input_file: Path to input .pptx file
        output_file: Path to output .pptx file
        spacing: Line spacing multiplier (e.g., 1.0, 1.5, 2.0) or Pt value as string
    """
    print(f"Loading presentation from {input_file}...")
    prs = Presentation(input_file)

    # Parse spacing value
    use_points = False
    if isinstance(spacing, str) and spacing.lower().endswith('pt'):
        use_points = True
        pt_value = float(spacing.lower().replace('pt', ''))
        print(f"Using exact point spacing: {pt_value} pt")
    else:
        spacing = float(spacing)
        print(f"Using line spacing multiplier: {spacing}")

    slide_count = 0
    shape_count = 0
    paragraph_count = 0

    # Iterate through all slides
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_count += 1

        # Iterate through all shapes in the slide
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                shape_count += 1
                text_frame = shape.text_frame

                # Iterate through all paragraphs in the text frame
                for paragraph in text_frame.paragraphs:
                    paragraph_count += 1

                    # Set line spacing
                    if use_points:
                        # Set exact point spacing
                        paragraph.line_spacing = Pt(pt_value)
                    else:
                        # Set relative spacing (1.0 = single, 1.5 = 1.5x, 2.0 = double)
                        paragraph.line_spacing = spacing

    # Save the modified presentation
    print(f"\nProcessed:")
    print(f"  {slide_count} slides")
    print(f"  {shape_count} text boxes")
    print(f"  {paragraph_count} paragraphs")
    print(f"\nSaving presentation to {output_file}...")
    prs.save(output_file)
    print("Done!")

def main():
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2]
    spacing = 1.0

    # Parse optional spacing argument
    i = 3
    while i < len(sys.argv):
        if sys.argv[i] == '--spacing' and i + 1 < len(sys.argv):
            spacing = sys.argv[i + 1]
            i += 2
        else:
            i += 1

    adjust_line_spacing(input_file, output_file, spacing)

if __name__ == '__main__':
    main()
